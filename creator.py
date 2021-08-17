import os
from pptx import Presentation, util, text
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX


from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE


def start():
    print("*********************************************")
    print("**Bem-vindo ao script de criação dos slides**")
    print("*********************************************")


def files_verification():
    print("Coloque a música que deseja na mesma pasta desse script. Tecle 'Enter' ao finalizar")
    status = input("")


def archive_list():
    global files, dir
    dir = os.getcwd()
    dir = dir + "/songs/"
    for root, dirs, files in os.walk(dir):
        songs = [files]
    return files


def select_song():
    global selected_file, complete_dir_file
    print("Selecione a música:")
    index = 0
    for i in files:
        print("({0}) - {1}".format(index, files[index]))
        index += 1
    selected_index = int(input(""))
    selected_file = files[selected_index]
    print("Você selecionou a música: {}".format(selected_file))
    complete_dir_file = dir + selected_file


def read_lyrics():
    lyric = open(complete_dir_file, 'r')
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = selected_file.replace(".txt", "").replace("-", " ").upper()
    i = 0
    par = []

    for line in lyric:
        line = line.strip('\n')
        if(line != ""):
            if(i == 0):
                paragraph = line
            else:
                paragraph = paragraph + "\n" + line
        else:
            par.append(paragraph)
            paragraph = ""

        i += 1
    par.append(paragraph)

    for ls in par:

        lyric_slide = prs.slide_layouts[6]
        slide = prs.slides.add_slide(lyric_slide)

        title = slide.shapes.add_textbox(left = Cm(0), top = Cm(5), width = Cm(20), height = Cm(10))

        #title.word_wrap = True

        #title.left = Cm(10)
        #title.top = Pt(0)
        #title.width = Pt(10)
        #title.height = Pt(7.5)

        lyr = title.text_frame
        p = lyr.add_paragraph()
        p.text = ls
        p.font.size = Pt(30)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.word_wrap = True
        #p.left = Cm(10)
        #p.top = Cm(0)
        #p.width = Cm(5)
        #p.height =Cm(7.5)


    pptx_name = selected_file + ".pptx"
    pptx_name = pptx_name.replace(".txt", "")
    prs.save(pptx_name)



start()
files_verification()
archive_list()
select_song()
read_lyrics()
